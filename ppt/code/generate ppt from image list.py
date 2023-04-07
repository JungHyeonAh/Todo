import os     # 파일이나 디렉토리 조작 가능, path를 얻거나 새로운 파일 direc 작성 가능.
import glob   # 파일들을 검색할 때 사용
import sys    # 파이썬 인터프리터 제어 방법 제공.
import json
import cv2

from pptx import Presentation   # ppt 파일 만들고 업데이트
from pptx.enum.shapes import MSO_SHAPE    # 도형 삽입
from pptx.enum.dml import MSO_COLOR_TYPE  # 색상 
from pptx.dml.color import RGBColor # 도형색을 수로 지정
from pptx.util import Inches, Pt    # Inches : 도형 크기, 여백 길이 지정 / Pt : 폰트 사이즈
from shutil import copyfile    # shutil : 폴더, 파일 복사


# Hospital = ("해맑은", "어튼", "누리꿈", "마디와", "미래", "조은", "한빛")
# Nerve = ("ISBP", "SCBP", "ICBP", "AxBP", "EbMN", "EbRN", "EbUN", "WrRN")
Hospital = ("a", "b", "o")
Nerve = ("A", "B", "O")


# root_path = '//10.50.29.23/cad/Nerve_2nd/Data_backup/'
root_path = 'C:\pptx\image'
nNerve = 3 # AxBP

def add_nerve_rectangle(slide, left, top, width, height, color, label) -> None:
    # 모서리가 둥근 사각형 도형 생성
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()   # 배경색 설정
    line = shape.line
    if color is 'RED':
        line.color.rgb = RGBColor(255, 0, 0)
    elif color is 'GREEN':
        line.color.rgb = RGBColor(0, 255, 0)
    elif color is 'YELLOW':
        line.color.rgb = RGBColor(250, 255, 0)

    txBox = slide.shapes.add_textbox(left, top-Pt(12), width, height)
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(8)
    return
    
def read_json(json_file_name) -> list:
    rect_cord = []
    with open(json_file_name, 'r') as f:
        annotation_content = json.loads(f.read())
        if "error_code" not in annotation_content:
            # update label id in the json file
            for data in annotation_content["data"]:
                if data["labels"][0]["name"] == 'Artery' or 'Nerve' in data["labels"][0]["name"]:
                    for shapes in data["shapes"]:
                        points = shapes["geometry"]["points"]
                        max_x = max_y = 0
                        min_x = min_y = 1
                        for cord in points:
                            if cord['x'] > max_x:
                                max_x = cord['x']
                            if cord['y'] > max_y:
                                max_y = cord['y']
                            if cord['x'] < min_x:
                                min_x = cord['x']
                            if cord['y'] < min_y:
                                min_y = cord['y']
                        width = max_x - min_x
                        height = max_y - min_y
                        left = min_x
                        top = min_y
                        cord = [left, top, width, height, data["labels"][0]['name']]
                        rect_cord.append(cord)
    return rect_cord

def copy_images_with_json() -> bool:
    for nHspt in [1]:# range(0,7):
        os.chdir(root_path + Hospital[nHspt])
        json_folder_list = sorted(glob.glob('*_json'))
        for json_path in json_folder_list:
            if Nerve[nNerve] in json_path:
                os.chdir(root_path + Hospital[nHspt] + '/' + json_path)
                json_list = sorted(glob.glob('*.json'))
                for json_name in json_list:
                    i = json_name[len(json_name)::-1].find('_')
                    image_name = json_name[:-i-1] + '.jpg'
                    try:
                        copyfile(root_path + Hospital[nHspt] + '/' + json_path[:len(json_path)-5] + '/' + image_name, root_path + 'Review/' + Hospital[nHspt] + '/' + json_path[:len(json_path)-5]+ '/' + image_name)
                    except:
                        print('copy error')
    return True



def main() -> None:

    #copy_images_with_json()

    for nHspt in [4, 5]:# range(0,7): 
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        title_slide_layout = prs.slide_layouts[0]
        os.chdir(root_path + 'Review/' + Hospital[nHspt])
        pjt_list = sorted(glob.glob('*'))
        for pjt_name in pjt_list:
            if Nerve[nNerve] in pjt_name and 'ppt' not in pjt_name:
                os.chdir(root_path + 'Review/' + Hospital[nHspt] + '/' + pjt_name)
                img_list = sorted(glob.glob('*.jpg'))
                json_path = root_path + Hospital[nHspt] + '/' + pjt_name + '_json/'
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                title.text = pjt_name

                for img_name in img_list:
                    slide = prs.slides.add_slide(blank_slide_layout)
                    left = top = width = height = Inches(0.05)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = pjt_name + '/' + img_name

                    left = top = Inches(0.4)
                    height = Inches(7)
                    pic = slide.shapes.add_picture(img_name, left, top, height=height)

                    img = cv2.imread(img_name)
                    height, width, layers = img.shape
                    asp_ratio = width/height

                    json_list = sorted(glob.glob(json_path + '*.json'))
                    img_name_short = img_name[:len(img_name)-4]
                    for json_name in json_list:
                        if img_name_short in json_name:
                            print(img_name_short)
                            nerve_rect = read_json(json_name)
                            for nRect in range(0, len(nerve_rect)):
                                add_nerve_rectangle(slide, left + nerve_rect[nRect][0]*Inches(7)*asp_ratio, top + nerve_rect[nRect][1]*Inches(7), nerve_rect[nRect][2]*Inches(7)*asp_ratio, nerve_rect[nRect][3]*Inches(7), 'YELLOW', nerve_rect[nRect][4])

        prs.save(root_path + 'Review/' + Hospital[nHspt] + '/' + Hospital[nHspt] + '_' + Nerve[nNerve] + '.pptx')
    return

if __name__ == "__main__":
    sys.exit(main() or 0)
